"""
extract_brand_icons.py — извлекает иконки из .pptx-шаблона (SberF1) в готовые
PNG-ассеты для генератора (assets/icons/icon_<name>_{dark,light}.png).

Иконки в шаблоне — векторные фигуры (custGeom) на слайдах «Набор иконок».
Скрипт:
  1. рендерит слайд на белом фоне (иконки тёмные),
  2. кластеризует штрихи в иконки (union-find по близости bbox),
  3. нарезает каждую иконку и собирает пронумерованный контактный лист,
  4. по выбранной карте {имя: индекс} сохраняет прозрачные PNG в двух
     вариантах (тёмный #111827 и светлый #FFFFFF) — они перекрашиваемые
     под любую подложку и одинаково рендерятся в Р7-Офис/PowerPoint/LibreOffice.

Зависимости: python-pptx, lxml, Pillow, LibreOffice (soffice), poppler (pdftoppm).

Пример:
    # 1) контактный лист — посмотреть индексы
    python -m pptx_generator.tools.extract_brand_icons \
        --template SberF1.pptx --slide 6 --contact-sheet contact.png
    # 2) извлечь выбранные (правишь CHOSEN в этом файле или передаёшь --map)
    python -m pptx_generator.tools.extract_brand_icons \
        --template SberF1.pptx --slide 6 --out assets/icons
"""
from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"a": A, "p": P}

# Карта по умолчанию: смысловое имя -> индекс кластера на контактном листе
# (для слайда 6 шаблона SberF1). Под другой слайд/шаблон — пересними лист.
CHOSEN = {
    "document": 14, "tag": 25, "rocket": 26, "growth": 32, "chart": 33,
    "integration": 35, "data": 43, "users": 49, "calendar": 52, "check": 54,
    "wifi": 59, "error": 60, "new": 61, "info": 62, "warning": 63,
    "list": 65, "security": 79,
}
DARK = (0x11, 0x18, 0x27)
LIGHT = (0xFF, 0xFF, 0xFF)


def _icon_boxes(slide_xml: Path):
    root = etree.parse(str(slide_xml)).getroot()
    boxes = []
    for sp in root.findall(".//p:sp", NS):
        if sp.find(".//a:custGeom", NS) is None:
            continue
        off = sp.find(".//a:xfrm/a:off", NS)
        ext = sp.find(".//a:xfrm/a:ext", NS)
        if off is None or ext is None:
            continue
        boxes.append((int(off.get("x")), int(off.get("y")),
                      int(ext.get("cx")), int(ext.get("cy"))))
    return boxes


def _cluster(boxes, pad=70000):
    def inflate(b):
        x, y, cx, cy = b
        return (x - pad, y - pad, x + cx + pad, y + cy + pad)

    R = [inflate(b) for b in boxes]
    par = list(range(len(boxes)))

    def find(i):
        while par[i] != i:
            par[i] = par[par[i]]
            i = par[i]
        return i

    def overlap(a, b):
        return not (a[2] < b[0] or b[2] < a[0] or a[3] < b[1] or b[3] < a[1])

    for i in range(len(R)):
        for j in range(i + 1, len(R)):
            if overlap(R[i], R[j]):
                par[find(i)] = find(j)
    groups = {}
    for i in range(len(boxes)):
        groups.setdefault(find(i), []).append(i)
    clusters = []
    for g in groups.values():
        xs = [boxes[i][0] for i in g]
        ys = [boxes[i][1] for i in g]
        xe = [boxes[i][0] + boxes[i][2] for i in g]
        ye = [boxes[i][1] + boxes[i][3] for i in g]
        clusters.append((min(xs), min(ys), max(xe), max(ye)))
    clusters.sort(key=lambda c: (round(c[1] / 150000), c[0]))
    return clusters


def _render_white(template: Path, slide_idx: int, workdir: Path) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    tmp = workdir / "clean.pptx"
    shutil.copy(template, tmp)
    prs = Presentation(str(tmp))
    lst = prs.slides._sldIdLst
    for i, sid in enumerate(list(lst)):
        if i != slide_idx:
            lst.remove(sid)
    f = prs.slides[0].background.fill
    f.solid()
    f.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    prs.save(str(tmp))
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                    "--outdir", str(workdir), str(tmp)], check=True,
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    subprocess.run(["pdftoppm", "-png", "-r", "200",
                    str(workdir / "clean.pdf"), str(workdir / "clean")],
                   check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return workdir / "clean-1.png"


def _make_variant(crop, rgb):
    from PIL import Image
    im = crop.convert("L")
    w, h = im.size
    px = im.load()
    out = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    op = out.load()
    for y in range(h):
        for x in range(w):
            a = 255 - px[x, y]
            op[x, y] = (rgb[0], rgb[1], rgb[2], 0 if a < 18 else a)
    bb = out.getbbox()
    if bb:
        out = out.crop(bb)
    s = int(max(out.size) * 1.16)
    canvas = Image.new("RGBA", (s, s), (0, 0, 0, 0))
    canvas.paste(out, ((s - out.width) // 2, (s - out.height) // 2), out)
    return canvas


def main():
    from PIL import Image, ImageDraw
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True, type=Path)
    ap.add_argument("--slide", type=int, default=6, help="номер слайда (1-индекс)")
    ap.add_argument("--out", type=Path, default=None, help="папка для PNG-иконок")
    ap.add_argument("--contact-sheet", type=Path, default=None)
    ap.add_argument("--map", type=Path, default=None, help="JSON {имя: индекс}")
    args = ap.parse_args()

    chosen = CHOSEN
    if args.map:
        chosen = json.loads(args.map.read_text(encoding="utf-8"))

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        # распаковать нужный slideN.xml
        import zipfile
        with zipfile.ZipFile(args.template) as z:
            name = f"ppt/slides/slide{args.slide}.xml"
            (td / "slide.xml").write_bytes(z.read(name))
        boxes = _icon_boxes(td / "slide.xml")
        clusters = _cluster(boxes)
        render = _render_white(args.template, args.slide - 1, td)
        img = Image.open(render).convert("RGB")
        W = img.size[0]
        sc = W / 12192000

        crops = []
        for (x0, y0, x1, y1) in clusters:
            px = (int(x0 * sc) - 6, int(y0 * sc) - 6, int(x1 * sc) + 6, int(y1 * sc) + 6)
            crops.append(img.crop(px))

        if args.contact_sheet:
            cell = 150
            cols = 8
            rows = (len(crops) + cols - 1) // cols
            sheet = Image.new("RGB", (cols * cell, rows * cell), (255, 255, 255))
            d = ImageDraw.Draw(sheet)
            for i, c in enumerate(crops):
                s = min(110 / c.width, 110 / c.height, 4)
                c2 = c.resize((int(c.width * s), int(c.height * s)))
                r, cc = divmod(i, cols)
                sheet.paste(c2, (cc * cell + (cell - c2.width) // 2, r * cell + 24))
                d.text((cc * cell + 5, r * cell + 4), f"#{i}", fill=(220, 0, 0))
            sheet.save(args.contact_sheet)
            print(f"contact sheet -> {args.contact_sheet} ({len(crops)} icons)")

        if args.out:
            os.makedirs(args.out, exist_ok=True)
            for name, idx in chosen.items():
                if idx >= len(crops):
                    print(f"skip {name}: index {idx} out of range")
                    continue
                _make_variant(crops[idx], DARK).save(args.out / f"icon_{name}_dark.png")
                _make_variant(crops[idx], LIGHT).save(args.out / f"icon_{name}_light.png")
            print(f"saved {len(chosen)} icons (x2 variants) -> {args.out}")


if __name__ == "__main__":
    main()
